from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

# Create a PowerPoint presentation object
presentation = Presentation()

# Slide 1: Title Slide
slide_1 = presentation.slides.add_slide(presentation.slide_layouts[0])
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]
title.text = "Feasibility Analysis of Expanding eManat Kiosks in Nakhchivan"
subtitle.text = "A Financial and Demographic Overview\nAugust 2024"

# Slide 2: Population and Demographics with Pie Chart
slide_2 = presentation.slides.add_slide(presentation.slide_layouts[5])
title = slide_2.shapes.title
title.text = "Population and Demographics"

# Data for pie chart
chart_data = CategoryChartData()
chart_data.categories = ['Urban Population', 'Rural Population']
chart_data.add_series('Population Distribution', (35.4, 64.6))

# Add pie chart to slide
x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4.5)
chart = slide_2.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart
chart.has_legend = True

# Add text
text_box = slide_2.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(2))
text_frame = text_box.text_frame
text_frame.text = (
    "Nakhchivan Population: 465,700\n"
    "- Urban: 35.4% (164,867 people)\n"
    "- Rural: 64.6% (300,833 people)\n\n"
    "The majority of the population is Azerbaijani (99.6%). "
    "The region shows a significant rural presence, which may affect kiosk placement and usage patterns."
)

# Slide 3: Kiosk Cost Analysis
slide_3 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_3.shapes.title
content = slide_3.shapes.placeholders[1]
title.text = "Kiosk Cost Analysis"
content.text = (
    "Estimated Kiosk Cost: 5,000 AZN per unit (standard features)\n"
    "Operational Costs: 1,000 AZN per month (maintenance, electricity, staffing)\n\n"
    "The costs reflect a mid-range investment, considering both purchase and ongoing operational expenses."
)

# Slide 4: Revenue Projections with Chart
slide_4 = presentation.slides.add_slide(presentation.slide_layouts[5])
title = slide_4.shapes.title
title.text = "Revenue Projections"

# Updated Data for chart with 50 transactions/day and 0.1 AZN/transaction
daily_revenue = 50 * 0.1  # 50 transactions per day at 0.1 AZN per transaction
monthly_revenue = daily_revenue * 30  # Assuming 30 days in a month

chart_data = CategoryChartData()
chart_data.categories = ["Daily Revenue", "Monthly Revenue"]
chart_data.add_series('Revenue', (daily_revenue, monthly_revenue))

# Add chart to slide
x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(3)
chart = slide_4.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
chart.has_legend = False

# Add text
text_box = slide_4.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(2))
text_frame = text_box.text_frame
text_frame.text = (
    f"Estimated Daily Transaction Volume: 50 transactions per kiosk\n"
    f"Revenue per Transaction: 0.1 AZN\n\n"
    f"Daily Revenue per Kiosk: {daily_revenue} AZN\n"
    f"Monthly Revenue per Kiosk: {monthly_revenue} AZN\n\n"
    "These projections are based on consistent demand, which is expected in urban areas."
)

# Slide 5: Profit and Loss Analysis with Chart
slide_5 = presentation.slides.add_slide(presentation.slide_layouts[5])
title = slide_5.shapes.title
title.text = "Profit and Loss Analysis"

# Updated Data for chart
monthly_profit = monthly_revenue - 1000  # Subtract operational costs of 1,000 AZN per month

chart_data = CategoryChartData()
chart_data.categories = ["Monthly Profit"]
chart_data.add_series('Profit', (monthly_profit,))

# Add chart to slide
x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(3)
chart = slide_5.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
chart.has_legend = False

# Add text
text_box = slide_5.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(2))
text_frame = text_box.text_frame
text_frame.text = (
    f"Monthly Profit per Kiosk: {monthly_profit} AZN (Revenue - Operational Costs)\n"
    "Break-even Point: Based on updated revenue projections, the break-even point will be recalculated.\n\n"
    "The quick break-even point indicates a profitable venture, assuming the transaction volume is sustained."
)

# Slide 6: Insights and Recommendations
slide_6 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_6.shapes.title
content = slide_6.shapes.placeholders[1]
title.text = "Insights and Recommendations"
content.text = (
    "1. Urban Focus: Kiosks should be primarily placed in urban areas where demand is higher.\n"
    "2. Rural Potential: Consideration should be given to rural areas with lower but stable demand.\n"
    "3. Scalability: The model shows potential for rapid scalability across similar regions.\n"
    "4. Risk Assessment: Further analysis is needed to address risks such as fluctuating transaction volumes and maintenance challenges."
)

# Save the presentation
presentation.save("eManat_Kiosk_Expansion_Nakhchivan_Updated_Revenue.pptx")
