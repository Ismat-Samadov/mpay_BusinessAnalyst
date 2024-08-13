from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt

# Corrected data from the text
data_2022 = {
    "Total Card Payments": 55980,  # in million manats
    "Cash Payments": 31787,  # in million manats
    "Cashless Payments": 24192,  # in million manats
    "POS Terminals": 5803 ,  # in million manats to million manats
    "E-commerce": 18348  # in million manats
}

# Insights to accompany the charts
insights = (
    "1. Dominance of Cash Payments: In 2022, cash payments accounted for 56.8% of the total "
    "card payments in Azerbaijan, highlighting a continued reliance on cash.\n\n"
    "2. Significant Shift to Cashless Payments: Despite the dominance of cash, 43.2% of the total "
    "payments were cashless, indicating a growing trend towards digital payment methods.\n\n"
    "3. E-commerce Leading Cashless Transactions: Within cashless payments, e-commerce "
    "transactions made up 76% of the total, demonstrating the increasing preference for online shopping "
    "among consumers. POS terminals accounted for 24%, but in absolute terms, they processed a larger "
    "amount due to the high volume of transactions.\n\n"
    "4. Digital Economy Growth: The significant volume of transactions through both POS terminals "
    "and e-commerce platforms reflects the expanding digital economy in Azerbaijan, with a clear shift "
    "towards cashless and online transactions."
)

# Create a PowerPoint presentation object
presentation = Presentation()

# Function to save plot as an image and add to slide
def add_plot_to_slide(slide, title):
    plt.savefig("plot.png", bbox_inches='tight')
    slide.shapes.title.text = title
    slide.shapes.add_picture("plot.png", Inches(1), Inches(1), width=Inches(8), height=Inches(4))
    plt.close()

# Slide 1: Cash vs Cashless Payments in 2022
plt.figure(figsize=(6, 4))
labels = ['Cash Payments', 'Cashless Payments']
sizes = [data_2022["Cash Payments"], data_2022["Cashless Payments"]]
colors = ['lightcoral', 'lightskyblue']
plt.pie(sizes, labels=labels, colors=colors, autopct='%1.1f%%', startangle=140)
plt.title('Cash vs Cashless Payments in 2022')

slide_1 = presentation.slides.add_slide(presentation.slide_layouts[5])
add_plot_to_slide(slide_1, "Cash vs Cashless Payments in 2022")

# Slide 2: POS Terminals vs E-commerce Payments in 2022
plt.figure(figsize=(6, 4))
labels = ['POS Terminals', 'E-commerce']
sizes = [data_2022["POS Terminals"], data_2022["E-commerce"]]
colors = ['gold', 'lightgreen']
plt.pie(sizes, labels=labels, colors=colors, autopct='%1.1f%%', startangle=140)
plt.title('POS Terminals vs E-commerce Payments in 2022')

slide_2 = presentation.slides.add_slide(presentation.slide_layouts[5])
add_plot_to_slide(slide_2, "POS Terminals vs E-commerce Payments in 2022")

# Insights Slide
slide_3 = presentation.slides.add_slide(presentation.slide_layouts[1])
slide_3.shapes.title.text = "Insights and Analysis"
text_box = slide_3.shapes.placeholders[1].text_frame
text_box.text = insights

# Save the presentation
presentation.save("Payment_Analysis_Azerbaijan_Fixed.pptx")
