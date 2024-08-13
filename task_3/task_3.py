from pptx import Presentation
from pptx.util import Inches
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns

# Load the Excel file
file_path = 'payment_data.xlsx'
data = pd.read_excel(file_path)

# Convert payment date to datetime format
data['payment date'] = pd.to_datetime(data['payment date'])

# Select the top 5 service categories based on total payment amount
top_service_categories = data.groupby('service category')['payment amount'].sum().sort_values(ascending=False).head(5).index

# Filter data to only include the top 5 service categories
data_top = data[data['service category'].isin(top_service_categories)]

# Divide the month into three parts: beginning (1-10), middle (11-20), and end (21-31)
data_top['month_part'] = pd.cut(data_top['payment date'].dt.day,
                                bins=[0, 10, 20, 31],
                                labels=['Beginning', 'Middle', 'End'])

# Calculate the average monthly spending per customer
monthly_spending = data_top.groupby([data_top['payment date'].dt.to_period('M'), 'customer ID'])['payment amount'].sum()
average_monthly_spending = monthly_spending.groupby('customer ID').mean()

# Identify below-average spenders
below_average_spenders = average_monthly_spending[average_monthly_spending < average_monthly_spending.mean()]
below_average_customers = below_average_spenders.index

# Filter data for below-average spenders
below_average_data = data_top[data_top['customer ID'].isin(below_average_customers)]

# Create a PowerPoint presentation
presentation = Presentation()

# 1. Customer Segmentation: Transaction Frequency
customer_activity = data_top['customer ID'].value_counts()
plt.figure(figsize=(8, 6))
customer_activity.plot(kind='hist', bins=20, color='orange')
plt.title('Customer Activity - Number of Transactions per Customer')
plt.xlabel('Number of Transactions')
plt.ylabel('Number of Customers')
plt.savefig('customer_activity.png', bbox_inches='tight')
plt.close()

slide_1 = presentation.slides.add_slide(presentation.slide_layouts[5])
slide_1.shapes.title.text = 'Customer Activity - Number of Transactions per Customer'
slide_1.shapes.add_picture('customer_activity.png', Inches(1), Inches(1), width=Inches(8), height=Inches(4))

# 2. Customer Segmentation: High Value vs Low Value Customers
high_value_customers = data_top.groupby('customer ID')['payment amount'].sum()
plt.figure(figsize=(8, 6))
high_value_customers.plot(kind='hist', bins=20, color='green')
plt.title('High Value vs Low Value Customers')
plt.xlabel('Total Payment Amount')
plt.ylabel('Number of Customers')
plt.savefig('high_value_customers.png', bbox_inches='tight')
plt.close()

slide_2 = presentation.slides.add_slide(presentation.slide_layouts[5])
slide_2.shapes.title.text = 'High Value vs Low Value Customers'
slide_2.shapes.add_picture('high_value_customers.png', Inches(1), Inches(1), width=Inches(8), height=Inches(4))

# 3. Payment Failure Analysis: Distribution of Failures by Service Category (Top 5)
failed_payments = data_top[data_top['payment status'] == 'failed']
failed_service_dist = failed_payments['service category'].value_counts()
plt.figure(figsize=(8, 6))
failed_service_dist.plot(kind='bar', color='red')
plt.title('Payment Failure Distribution by Service Category (Top 5)')
plt.ylabel('Number of Failed Payments')
plt.savefig('failed_service_dist.png', bbox_inches='tight')
plt.close()

slide_3 = presentation.slides.add_slide(presentation.slide_layouts[5])
slide_3.shapes.title.text = 'Payment Failure Distribution by Service Category (Top 5)'
slide_3.shapes.add_picture('failed_service_dist.png', Inches(1), Inches(1), width=Inches(8), height=Inches(4))

# 4. Monthly Trends with Seasonal Decomposition (Top 5)
monthly_payments = data_top.groupby(data_top['payment date'].dt.to_period('M'))['payment amount'].sum()
plt.figure(figsize=(8, 6))
monthly_payments.plot(kind='line', marker='o', color='purple')
plt.title('Monthly Payment Trends with Seasonal Decomposition (Top 5)')
plt.ylabel('Total Payment Amount (Manats)')
plt.xlabel('Month')
plt.savefig('monthly_payments.png', bbox_inches='tight')
plt.close()

slide_4 = presentation.slides.add_slide(presentation.slide_layouts[5])
slide_4.shapes.title.text = 'Monthly Payment Trends with Seasonal Decomposition (Top 5)'
slide_4.shapes.add_picture('monthly_payments.png', Inches(1), Inches(1), width=Inches(8), height=Inches(4))

# 5. Spending Patterns Across the Month (Top 5)
spending_patterns = data_top.groupby(['month_part', 'service category'])['payment amount'].sum().unstack().fillna(0)
plt.figure(figsize=(8, 6))
spending_patterns.plot(kind='bar', stacked=True)
plt.title('Spending Patterns Across the Month (Top 5)')
plt.ylabel('Total Payment Amount (Manats)')
plt.xlabel('Part of the Month')
plt.savefig('spending_patterns.png', bbox_inches='tight')
plt.close()

slide_5 = presentation.slides.add_slide(presentation.slide_layouts[5])
slide_5.shapes.title.text = 'Spending Patterns Across the Month (Top 5)'
slide_5.shapes.add_picture('spending_patterns.png', Inches(1), Inches(1), width=Inches(8), height=Inches(4))

# 6. Correlation Analysis Between Part of the Month and Payment Amount
correlation_data = data_top[['month_part', 'payment amount']].copy()
correlation_data['month_part'] = correlation_data['month_part'].apply(lambda x: {'Beginning': 1, 'Middle': 2, 'End': 3}[x])
correlation_matrix = correlation_data.corr()

plt.figure(figsize=(8, 6))
sns.heatmap(correlation_matrix, annot=True, cmap='coolwarm', center=0)
plt.title('Correlation Between Part of the Month and Payment Amount')
plt.savefig('correlation_matrix.png', bbox_inches='tight')
plt.close()

slide_6 = presentation.slides.add_slide(presentation.slide_layouts[5])
slide_6.shapes.title.text = 'Correlation Between Part of the Month and Payment Amount'
slide_6.shapes.add_picture('correlation_matrix.png', Inches(1), Inches(1), width=Inches(8), height=Inches(4))

# 7. Analysis of Below-Average Spenders
below_avg_spending_patterns = below_average_data.groupby(['month_part', 'service category'])['payment amount'].sum().unstack().fillna(0)
plt.figure(figsize=(8, 6))
below_avg_spending_patterns.plot(kind='bar', stacked=True, color=sns.color_palette('muted'))
plt.title('Spending Patterns of Below-Average Spenders')
plt.ylabel('Total Payment Amount (Manats)')
plt.xlabel('Part of the Month')
plt.savefig('below_avg_spending_patterns.png', bbox_inches='tight')
plt.close()

slide_7 = presentation.slides.add_slide(presentation.slide_layouts[5])
slide_7.shapes.title.text = 'Spending Patterns of Below-Average Spenders'
slide_7.shapes.add_picture('below_avg_spending_patterns.png', Inches(1), Inches(1), width=Inches(8), height=Inches(4))

# Save the presentation
presentation.save('payment_data_top5_desc_correlation.pptx')

print("Analysis saved to 'payment_data_top5_desc_correlation.pptx'")

